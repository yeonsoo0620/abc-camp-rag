#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Nordic Modern 스타일 도서 기획 프레젠테이션 생성 스크립트
예스24 IT/모바일 베스트셀러 데이터 기반 신규 도서 기획
"""

import csv
import os
import re
from collections import Counter, defaultdict

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData

# ── Nordic Modern Color Palette ──
SLATE_DARK   = RGBColor(0x2D, 0x3A, 0x4A)   # 2D3A4A - primary dark
SLATE_MID    = RGBColor(0x4A, 0x5D, 0x6E)   # 4A5D6E - secondary
WARM_WHITE   = RGBColor(0xF7, 0xF5, 0xF0)   # F7F5F0 - background
CREAM        = RGBColor(0xF0, 0xEB, 0xE1)   # F0EBE1 - light bg
WOOD_BROWN   = RGBColor(0xC4, 0xA8, 0x82)   # C4A8282 - accent
DEEP_TEAL    = RGBColor(0x3D, 0x7C, 0x88)   # 3D7C88 - teal accent
MUTED_GREEN  = RGBColor(0x7B, 0xA0, 0x8A)   # 7BA08A - sage
LIGHT_GRAY   = RGBColor(0xD6, 0xD1, 0xC7)   # D6D1C7 - border
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)   # FFFFFF
BLACK        = RGBColor(0x1A, 0x1A, 0x1A)   # 1A1A1A
MID_GRAY     = RGBColor(0x8A, 0x8A, 0x8A)   # 8A8A8A
CORAL        = RGBColor(0xE8, 0x6F, 0x51)   # E86F51 - warm accent
DARK_TEAL    = RGBColor(0x2C, 0x5F, 0x6A)   # 2C5F6A - deep teal
SOFT_BLUE    = RGBColor(0x8A, 0xBD, 0xC5)   # 8ABDC5 - soft blue

# Slide dimensions (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_BODY = "Calibri"
FONT_TITLE = "Georgia"

# ── Helper Functions ──

def add_bg(slide, color):
    """Set solid background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, x, y, w, h, fill_color):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, x, y, w, h, text, font_size=14, color=SLATE_DARK,
                 bold=False, italic=False, align=PP_ALIGN.LEFT, font_name=FONT_BODY,
                 valign=MSO_ANCHOR.TOP):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    tf.auto_size = None
    tf.vertical_anchor = valign
    return txBox

def add_multiline_text(slide, x, y, w, h, lines, font_size=14, color=SLATE_DARK,
                       bold=False, line_spacing=1.2, font_name=FONT_BODY, align=PP_ALIGN.LEFT):
    """Add text box with multiple lines."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(font_size * 0.3)
        if isinstance(line, tuple):
            text, opts = line
        else:
            text = line
            opts = {}
        run = p.add_run()
        run.text = text
        run.font.size = Pt(opts.get('size', font_size))
        run.font.color.rgb = opts.get('color', color)
        run.font.bold = opts.get('bold', bold)
        run.font.italic = opts.get('italic', False)
        run.font.name = opts.get('font', font_name)
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_circle(slide, x, y, size, fill_color):
    """Add a circle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def set_notes(slide, text):
    """Set speaker notes for a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ── Data Analysis Functions ──

def parse_sales_index(val):
    """Parse sales index value from string."""
    if not val or val.strip() == '':
        return 0
    val = val.strip().replace(',', '').replace('"', '')
    try:
        return int(val)
    except:
        return 0

def parse_price(val):
    """Parse price value from string."""
    if not val or val.strip() == '':
        return 0
    val = val.strip().replace(',', '').replace('"', '').replace('원', '')
    try:
        return int(val)
    except:
        return 0

def load_data(csv_path):
    """Load and parse the CSV data."""
    books = []
    with open(csv_path, 'r', encoding='utf-8-sig') as f:
        reader = csv.DictReader(f)
        for row in reader:
            book = {
                'rank': row.get('순위', ''),
                'title': row.get('도서명', ''),
                'author': row.get('저자', ''),
                'publisher': row.get('출판사', ''),
                'pub_date': row.get('출판일', ''),
                'sale_price': parse_price(row.get('판매가', '')),
                'orig_price': parse_price(row.get('정가', '')),
                'discount': row.get('할인율', ''),
                'sales_index': parse_sales_index(row.get('판매지수', ''))
            }
            books.append(book)
    return books

def categorize_book(title):
    """Categorize books by main topic."""
    title_lower = title.lower()
    if any(w in title_lower for w in ['claude', '클로드', 'code', '코워크']):
        return 'Claude/에이전틱'
    if any(w in title_lower for w in ['gpt', '챗지피티', '챗gpt', 'chatgpt']):
        return 'ChatGPT'
    if any(w in title_lower for w in ['gemini', '제미나이']):
        return 'Gemini'
    if any(w in title_lower for w in ['vibe', '바이브', '코딩']):
        return '바이브 코딩'
    if any(w in title_lower for w in ['교육', '교사', '수업', '에듀테크', '학교']):
        return '교육/에듀테크'
    if any(w in title_lower for w in ['ai', '인공지능', '생성형', '생성형 ai']):
        return 'AI/인공지능'
    if any(w in title_lower for w in ['영상', '유튜브', '숏츠', '쇼츠', '릴스']):
        return '영상/콘텐츠'
    if any(w in title_lower for w in ['엑셀', '노션', '피그마', '캔바', '디자인']):
        return '업무 도구'
    if any(w in title_lower for w in ['코딩', '파이썬', '프로그래밍', '언어']):
        return '코딩/프로그래밍'
    return '기타'


# ── Slide Generation Functions ──

def slide_01_title(prs):
    """Slide 1: Title Slide - Nordic dark background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, SLATE_DARK)

    # Decorative left bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, WOOD_BROWN)

    # Decorative circle elements
    add_circle(slide, Inches(10.5), Inches(0.8), Inches(1.8), SLATE_MID)
    add_circle(slide, Inches(11.2), Inches(2.0), Inches(1.2), DEEP_TEAL)

    # Title text
    add_text_box(slide, Inches(1.2), Inches(1.5), Inches(8), Inches(1.2),
                 "신규 도서 기획 제안서", font_size=44, color=WHITE,
                 bold=True, font_name=FONT_TITLE)

    # Subtitle
    add_text_box(slide, Inches(1.2), Inches(2.8), Inches(8), Inches(0.8),
                 "예스24 IT/모바일 베스트셀러 시장 분석 기반",
                 font_size=20, color=SOFT_BLUE, font_name=FONT_BODY)

    # Divider line
    add_rect(slide, Inches(1.2), Inches(3.8), Inches(3), Inches(0.04), WOOD_BROWN)

    # Date and info
    add_text_box(slide, Inches(1.2), Inches(4.2), Inches(6), Inches(0.5),
                 "2026년 7월  |  시장 조사 및 기획 보고서",
                 font_size=14, color=LIGHT_GRAY, font_name=FONT_BODY)

    # Bottom decorative bar
    add_rect(slide, Inches(0), Inches(6.8), SLIDE_W, Inches(0.7), SLATE_MID)
    add_text_box(slide, Inches(1.2), Inches(6.85), Inches(10), Inches(0.5),
                 "ABC Publishing  |  IT & Technology Division",
                 font_size=12, color=LIGHT_GRAY, font_name=FONT_BODY)

    # Speaker notes
    set_notes(slide, "안녕하세요. 오늘은 예스24 IT/모바일 베스트셀러 데이터를 바탕으로 "
              "신규 도서 기획안을 제안드리겠습니다. 2026년 현재 IT 도서 시장은 AI 관련 도서가 "
              "압도적인 인기를 끌고 있으며, 이 데이터를 분석하여 최적의 도서 기획 방향을 "
              "제시하겠습니다.")


def slide_02_overview(prs):
    """Slide 2: Project Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WARM_WHITE)

    # Top dark bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), SLATE_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                 "01  프로젝트 개요", font_size=32, color=WHITE,
                 bold=True, font_name=FONT_TITLE)

    # Left section - Objective
    add_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.2), WHITE)
    add_rect(slide, Inches(0.8), Inches(1.8), Inches(0.08), Inches(2.2), DEEP_TEAL)

    add_text_box(slide, Inches(1.2), Inches(1.9), Inches(4.8), Inches(0.5),
                 "프로젝트 목적", font_size=18, color=DEEP_TEAL, bold=True)

    add_multiline_text(slide, Inches(1.2), Inches(2.4), Inches(4.8), Inches(1.4),
        [
            "예스24 IT/모바일 종합 베스트셀러 250권 데이터 분석",
            "2025-2026년 IT 도서 시장 트렌드 파악",
            "차세대 베스트셀러 도서 기획 방향 도출",
            "타겟 독자층 및 마케팅 전략 수립",
        ], font_size=13, color=SLATE_MID, line_spacing=1.5)

    # Right section - Period
    add_rect(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(2.2), WHITE)
    add_rect(slide, Inches(7.0), Inches(1.8), Inches(0.08), Inches(2.2), WOOD_BROWN)

    add_text_box(slide, Inches(7.4), Inches(1.9), Inches(4.8), Inches(0.5),
                 "조사 기간 및 범위", font_size=18, color=WOOD_BROWN, bold=True)

    add_multiline_text(slide, Inches(7.4), Inches(2.4), Inches(4.8), Inches(1.4),
        [
            "조사 기간: 2025년 ~ 2026년 7월",
            "데이터 소스: 예스24 IT/모바일 베스트셀러",
            "수집 규모: 250권 종합 베스트셀러",
            "분석 항목: 가격, 판매지수, 카테고리, 출판사",
        ], font_size=13, color=SLATE_MID, line_spacing=1.5)

    # Bottom stats row
    stats = [
        ("250+", "분석 도서 수"),
        ("60+", "출판사"),
        ("80%", "AI 관련 도서 비중"),
        ("2.5만", "최고 판매지수"),
    ]

    for i, (num, label) in enumerate(stats):
        x = Inches(0.8 + i * 3.1)
        y = Inches(4.5)
        add_rect(slide, x, y, Inches(2.7), Inches(2.2), WHITE)
        add_circle(slide, x + Inches(0.9), y + Inches(0.3), Inches(0.9), CREAM)
        add_text_box(slide, x, y + Inches(0.35), Inches(2.7), Inches(0.8),
                     num, font_size=36, color=DEEP_TEAL, bold=True,
                     align=PP_ALIGN.CENTER, font_name=FONT_TITLE)
        add_text_box(slide, x, y + Inches(1.3), Inches(2.7), Inches(0.5),
                     label, font_size=13, color=SLATE_MID, align=PP_ALIGN.CENTER)

    set_notes(slide, "프로젝트 개요입니다. 이번 기획은 예스24 IT/모바일 종합 베스트셀러 "
              "250권의 데이터를 체계적으로 분석하여, 시장에서 성공할 수 있는 신규 도서를 "
              "기획하는 것이 목표입니다. 2025년부터 2026년 7월까지의 데이터를 확보하였으며, "
              "가격, 판매지수, 카테고리, 출판사 등 다양한 변수를 분석했습니다.")


def slide_03_market_overview(prs):
    """Slide 3: Market Overview with key stats"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WARM_WHITE)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), SLATE_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                 "02  IT 도서 시장 현황", font_size=32, color=WHITE,
                 bold=True, font_name=FONT_TITLE)

    # Market insight cards
    cards = [
        ("AI 도서 시장 폭발", "2026년 IT 베스트셀러의 80% 이상이\nAI/생성형 AI 관련 도서로, 시장이\n급격히 AI 중심으로 재편되고 있습니다.", DEEP_TEAL),
        ("바이브 코딩 열풍", "코딩을 모르는 일반인도 앱을 만들 수 있는\n바이브 코딩이 새로운 트렌드로 부상하며\n관련 도서 수요가 급증하고 있습니다.", WOOD_BROWN),
        ("교육 시장 확대", "교사 대상 에듀테크 도서가 지속적으로\n베스트셀러에 등장하며, 교육 시장이\n주요 타겟으로 떠오르고 있습니다.", MUTED_GREEN),
    ]

    for i, (title, desc, accent) in enumerate(cards):
        x = Inches(0.8 + i * 4.1)
        y = Inches(1.8)
        add_rect(slide, x, y, Inches(3.7), Inches(3.0), WHITE)
        add_rect(slide, x, y, Inches(3.7), Inches(0.08), accent)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.1), Inches(0.5),
                     title, font_size=18, color=accent, bold=True, font_name=FONT_TITLE)
        add_multiline_text(slide, x + Inches(0.3), y + Inches(0.9), Inches(3.1), Inches(1.8),
                           desc.split('\n'), font_size=12, color=SLATE_MID, line_spacing=1.6)

    # Key market data
    add_rect(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.8), SLATE_DARK)
    add_text_box(slide, Inches(1.2), Inches(5.4), Inches(4), Inches(0.5),
                 "핵심 시장 데이터", font_size=16, color=WOOD_BROWN, bold=True)

    market_data = [
        ("평균 판매가", "22,100원", "할인 적용 후 평균"),
        ("평균 할인율", "10%", "대부분 정가 대비 10% 할인"),
        ("최고 판매지수", "82,788", "AI 수업 활용 가이드"),
        ("신간 비중", "70%+", "2026년 출간 도서"),
    ]

    for i, (label, value, note) in enumerate(market_data):
        x = Inches(1.2 + i * 2.9)
        add_text_box(slide, x, Inches(5.9), Inches(2.5), Inches(0.3),
                     label, font_size=11, color=LIGHT_GRAY)
        add_text_box(slide, x, Inches(6.2), Inches(2.5), Inches(0.4),
                     value, font_size=22, color=WHITE, bold=True, font_name=FONT_TITLE)
        add_text_box(slide, x, Inches(6.7), Inches(2.5), Inches(0.3),
                     note, font_size=10, color=LIGHT_GRAY, italic=True)

    set_notes(slide, "IT 도서 시장 현황을 살펴보면, 2026년 현재 AI 관련 도서가 베스트셀러의 "
              "80% 이상을 차지하고 있습니다. 특히 바이브 코딩 열풍은 코딩을 모르는 일반인도 "
              "앱을 만들 수 있는 새로운 트렌드로, 관련 도서 수요가 급증하고 있습니다. "
              "교육 시장도 주요 타겟으로 부상하고 있으며, 평균 판매가는 약 22,100원, "
              "대부분 10% 할인율을 적용하고 있습니다.")


def slide_04_top10(prs, books):
    """Slide 4: Top 10 Bestsellers"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WARM_WHITE)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), SLATE_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                 "03  TOP 10 베스트셀러", font_size=32, color=WHITE,
                 bold=True, font_name=FONT_TITLE)

    # Sort by sales index
    sorted_books = sorted(books, key=lambda x: x['sales_index'], reverse=True)[:10]

    # Table header
    header_y = Inches(1.6)
    add_rect(slide, Inches(0.8), header_y, Inches(11.7), Inches(0.5), SLATE_MID)
    headers = [("순위", 0.8, 0.7), ("도서명", 1.5, 4.5), ("저자", 6.0, 1.8),
               ("출판사", 7.8, 1.5), ("판매지수", 9.3, 1.5), ("가격", 10.8, 1.2)]
    for text, x, w in headers:
        add_text_box(slide, Inches(x), header_y + Inches(0.08), Inches(w), Inches(0.35),
                     text, font_size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Table rows
    for i, book in enumerate(sorted_books):
        row_y = Inches(2.15 + i * 0.48)
        bg_color = WHITE if i % 2 == 0 else CREAM
        add_rect(slide, Inches(0.8), row_y, Inches(11.7), Inches(0.45), bg_color)

        # Rank
        rank_circle = add_circle(slide, Inches(1.0), row_y + Inches(0.05), Inches(0.35), DEEP_TEAL)
        add_text_box(slide, Inches(1.0), row_y + Inches(0.05), Inches(0.35), Inches(0.35),
                     str(i+1), font_size=11, color=WHITE, bold=True,
                     align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

        # Title (truncate if too long)
        title = book['title']
        if len(title) > 30:
            title = title[:28] + '...'
        add_text_box(slide, Inches(1.5), row_y + Inches(0.08), Inches(4.3), Inches(0.35),
                     title, font_size=11, color=SLATE_DARK, bold=True)

        # Author
        author = book['author']
        if len(author) > 15:
            author = author[:13] + '...'
        add_text_box(slide, Inches(6.0), row_y + Inches(0.08), Inches(1.7), Inches(0.35),
                     author, font_size=10, color=SLATE_MID)

        # Publisher
        add_text_box(slide, Inches(7.8), row_y + Inches(0.08), Inches(1.4), Inches(0.35),
                     book['publisher'], font_size=10, color=SLATE_MID)

        # Sales Index
        add_text_box(slide, Inches(9.3), row_y + Inches(0.08), Inches(1.4), Inches(0.35),
                     f"{book['sales_index']:,}", font_size=11, color=DEEP_TEAL, bold=True,
                     align=PP_ALIGN.RIGHT)

        # Price
        add_text_box(slide, Inches(10.8), row_y + Inches(0.08), Inches(1.2), Inches(0.35),
                     f"₩{book['sale_price']:,}", font_size=10, color=SLATE_MID,
                     align=PP_ALIGN.RIGHT)

    set_notes(slide, "TOP 10 베스트셀러를 살펴보면, '요즘 교사를 위한 AI 수업 활용 가이드'가 "
              "판매지수 82,788로 1위를 차지하고 있으며, '혼자 공부하는 바이브 코딩 with "
              "클로드 코드', '이게 되네? 제미나이 완전 미친 활용법 81제' 등이 뒤를 잇고 있습니다. "
              "특히 AI 관련 도서가 전체 10위권을 독식하고 있으며, 교육용 도서와 실전 활용서가 "
              "인기를 끌고 있습니다.")


def slide_05_publisher_analysis(prs, books):
    """Slide 5: Publisher Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WARM_WHITE)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), SLATE_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                 "04  출판사별 시장 점유율", font_size=32, color=WHITE,
                 bold=True, font_name=FONT_TITLE)

    # Publisher count
    pub_counter = Counter(b['publisher'] for b in books)
    top_pub = pub_counter.most_common(8)

    # Left: Bar chart (manual)
    chart_left = Inches(0.8)
    chart_top = Inches(1.8)
    chart_w = Inches(6.5)
    chart_h = Inches(5.0)

    add_rect(slide, chart_left, chart_top, chart_w, chart_h, WHITE)

    max_count = top_pub[0][1] if top_pub else 1
    bar_h = Inches(0.45)
    bar_gap = Inches(0.15)

    for i, (pub, count) in enumerate(top_pub[:7]):
        y = chart_top + Inches(0.3) + i * (bar_h + bar_gap)
        # Label
        add_text_box(slide, chart_left + Inches(0.2), y, Inches(2.0), bar_h,
                     pub, font_size=11, color=SLATE_DARK, valign=MSO_ANCHOR.MIDDLE)
        # Bar
        bar_width = (count / max_count) * Inches(3.5)
        colors = [DEEP_TEAL, WOOD_BROWN, MUTED_GREEN, CORAL, DARK_TEAL, SOFT_BLUE, SLATE_MID]
        add_rect(slide, chart_left + Inches(2.3), y + Inches(0.05),
                 int(bar_width), Inches(0.35), colors[i % len(colors)])
        # Count label
        add_text_box(slide, chart_left + Inches(2.3) + int(bar_width) + Inches(0.1), y,
                     Inches(0.8), bar_h, f"{count}권", font_size=11, color=SLATE_MID,
                     valign=MSO_ANCHOR.MIDDLE)

    # Right: Publisher insight cards
    right_x = Inches(7.8)
    add_rect(slide, right_x, Inches(1.8), Inches(4.7), Inches(5.0), WHITE)

    add_text_box(slide, right_x + Inches(0.3), Inches(2.0), Inches(4.1), Inches(0.5),
                 "출판사 인사이트", font_size=16, color=DEEP_TEAL, bold=True, font_name=FONT_TITLE)

    insights = [
        ("한빛미디어", "가장 많은 베스트셀러 보유. AI/코딩 도서 전문.", DEEP_TEAL),
        ("골든래빗", "실전 활용서 강세. '바로바로', '이게 되네?' 시리즈 성공.", WOOD_BROWN),
        ("이지스퍼블리싱", "'된다!' 시리즈로 브랜드 인지도 확보.", MUTED_GREEN),
        ("앤써북", "교육 시장 집중. 에듀테크 도서 강세.", CORAL),
        ("길벗", "기술 서적 전반. 신규 AI 도서 확대 중.", DARK_TEAL),
    ]

    for i, (pub, desc, accent) in enumerate(insights):
        y = Inches(2.7) + i * Inches(0.85)
        add_rect(slide, right_x + Inches(0.3), y, Inches(0.06), Inches(0.65), accent)
        add_text_box(slide, right_x + Inches(0.6), y, Inches(3.8), Inches(0.3),
                     pub, font_size=12, color=accent, bold=True)
        add_text_box(slide, right_x + Inches(0.6), y + Inches(0.3), Inches(3.8), Inches(0.35),
                     desc, font_size=10, color=SLATE_MID)

    set_notes(slide, "출판사별 시장 점유율을 분석하면, 한빛미디어가 가장 많은 베스트셀러를 "
              "보유하고 있으며, 골든래빗은 '바로바로', '이게 되네?' 시리즈로 실전 활용서 시장에서 "
              "강세를 보이고 있습니다. 이지스퍼블리싱은 '된다!' 시리즈로 브랜드 인지도를 확보했으며, "
              "앤써북은 교육 시장에 집중하고 있습니다. 신규 도서 기획 시 이 출판사들의 강점을 "
              "참고해야 합니다.")


def slide_06_price_analysis(prs, books):
    """Slide 6: Price Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WARM_WHITE)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), SLATE_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                 "05  가격 분석 및 포지셔닝", font_size=32, color=WHITE,
                 bold=True, font_name=FONT_TITLE)

    # Price distribution
    prices = [b['sale_price'] for b in books if b['sale_price'] > 0]
    avg_price = sum(prices) / len(prices) if prices else 0

    # Price ranges
    ranges = [
        ("15,000원 이하", len([p for p in prices if p <= 15000]), LIGHT_GRAY),
        ("15,000~20,000원", len([p for p in prices if 15000 < p <= 20000]), SOFT_BLUE),
        ("20,000~25,000원", len([p for p in prices if 20000 < p <= 25000]), DEEP_TEAL),
        ("25,000~30,000원", len([p for p in prices if 25000 < p <= 30000]), WOOD_BROWN),
        ("30,000원 초과", len([p for p in prices if p > 30000]), CORAL),
    ]

    max_range = max(r[1] for r in ranges) if ranges else 1

    # Price distribution chart
    add_rect(slide, Inches(0.8), Inches(1.8), Inches(7.0), Inches(3.5), WHITE)
    add_text_box(slide, Inches(1.1), Inches(1.9), Inches(4), Inches(0.5),
                 "가격대별 도서 분포", font_size=16, color=DEEP_TEAL, bold=True, font_name=FONT_TITLE)

    for i, (label, count, color) in enumerate(ranges):
        y = Inches(2.6) + i * Inches(0.55)
        add_text_box(slide, Inches(1.1), y, Inches(2.0), Inches(0.4),
                     label, font_size=11, color=SLATE_DARK)
        bar_w = (count / max_range) * Inches(3.5) if max_range > 0 else Inches(0)
        add_rect(slide, Inches(3.3), y + Inches(0.05), int(bar_w), Inches(0.3), color)
        add_text_box(slide, Inches(3.3) + int(bar_w) + Inches(0.1), y,
                     Inches(0.8), Inches(0.4), f"{count}권", font_size=11, color=SLATE_MID)

    # Right: Price insight
    add_rect(slide, Inches(8.2), Inches(1.8), Inches(4.3), Inches(3.5), WHITE)

    add_text_box(slide, Inches(8.5), Inches(2.0), Inches(3.7), Inches(0.5),
                 "가격 인사이트", font_size=16, color=WOOD_BROWN, bold=True, font_name=FONT_TITLE)

    price_insights = [
        f"평균 판매가: ₩{int(avg_price):,}",
        "가장 많은 가격대: 20,000~25,000원",
        "대부분 10% 할인 적용",
        "실전 서적은 25,000~30,000원대",
        "입문서는 18,000~22,000원대",
    ]

    for i, insight in enumerate(price_insights):
        y = Inches(2.7) + i * Inches(0.5)
        add_circle(slide, Inches(8.5), y + Inches(0.08), Inches(0.15), DEEP_TEAL)
        add_text_box(slide, Inches(8.8), y, Inches(3.4), Inches(0.4),
                     insight, font_size=12, color=SLATE_DARK)

    # Bottom recommendation
    add_rect(slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.5), SLATE_DARK)
    add_text_box(slide, Inches(1.2), Inches(5.75), Inches(3), Inches(0.4),
                 "가격 전략 제안", font_size=16, color=WOOD_BROWN, bold=True, font_name=FONT_TITLE)

    recs = [
        "입문서: 19,800원 (정가 22,000원, 10% 할인) - 대중성 확보",
        "실전서: 27,000원 (정가 30,000원, 10% 할인) - 전문성 반영",
        "올인원: 28,800~32,400원 - 프리미엄 포지셔닝",
    ]
    for i, rec in enumerate(recs):
        add_text_box(slide, Inches(1.2 + i * 3.9), Inches(6.2), Inches(3.6), Inches(0.6),
                     rec, font_size=11, color=LIGHT_GRAY)

    set_notes(slide, "가격 분석 결과, 20,000~25,000원대 도서가 가장 많으며, 평균 판매가는 "
              "약 22,100원입니다. 대부분의 도서가 10% 할인율을 적용하고 있습니다. "
              "입문서는 19,800원, 실전서는 27,000원, 올인원 서적은 28,800~32,400원대로 "
              "포지셔닝하는 것이 적절할 것으로 판단됩니다.")


def slide_07_category_trend_ai(prs, books):
    """Slide 7: AI/LLM Category Trends"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WARM_WHITE)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), SLATE_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                 "06  카테고리 트렌드: AI/LLM", font_size=32, color=WHITE,
                 bold=True, font_name=FONT_TITLE)

    # AI sub-categories
    ai_keywords = {
        'Claude/에이전틱': ['claude', '클로드', '코워크', '하네스', '에이전틱', '오픈클로', 'mcp'],
        'ChatGPT': ['gpt', '챗지피티', '챗gpt', 'chatgpt'],
        'Gemini': ['gemini', '제미나이'],
        '바이브 코딩': ['vibe', '바이브', '코딩'],
        '교육/에듀테크': ['교육', '교사', '수업', '에듀테크', '학교'],
    }

    ai_counts = {}
    for cat, keywords in ai_keywords.items():
        count = 0
        for b in books:
            if any(kw in b['title'].lower() for kw in keywords):
                count += 1
        ai_counts[cat] = count

    # Left: Category breakdown
    add_rect(slide, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), WHITE)
    add_text_box(slide, Inches(1.1), Inches(1.9), Inches(5), Inches(0.5),
                 "AI 관련 카테고리별 도서 수", font_size=16, color=DEEP_TEAL,
                 bold=True, font_name=FONT_TITLE)

    max_ai = max(ai_counts.values()) if ai_counts else 1
    colors = [DEEP_TEAL, WOOD_BROWN, MUTED_GREEN, CORAL, DARK_TEAL]

    for i, (cat, count) in enumerate(sorted(ai_counts.items(), key=lambda x: x[1], reverse=True)):
        y = Inches(2.7) + i * Inches(0.8)
        add_text_box(slide, Inches(1.1), y, Inches(2.2), Inches(0.4),
                     cat, font_size=12, color=SLATE_DARK, bold=True)
        bar_w = (count / max_ai) * Inches(3.0) if max_ai > 0 else Inches(0)
        add_rect(slide, Inches(3.4), y + Inches(0.05), int(bar_w), Inches(0.35),
                 colors[i % len(colors)])
        add_text_box(slide, Inches(3.4) + int(bar_w) + Inches(0.1), y,
                     Inches(0.8), Inches(0.4), f"{count}권", font_size=12,
                     color=SLATE_MID, bold=True)

    # Right: Trend insights
    add_rect(slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(5.0), WHITE)

    add_text_box(slide, Inches(7.5), Inches(2.0), Inches(4.7), Inches(0.5),
                 "AI 도서 시장 트렌드", font_size=16, color=WOOD_BROWN,
                 bold=True, font_name=FONT_TITLE)

    trends = [
        ("Claude Code 폭발적 성장", "클로드 코드 관련 도서가 급증하며\n코워크, 하네스 등 생태계 확장", DEEP_TEAL),
        ("Gemini 활용서 인기", "제미나이 완전 미친 활용법 등\n실전 활용서가 높은 판매지수 기록", WOOD_BROWN),
        ("교육 시장 AI 도입", "교사 대상 AI 수업 활용 가이드가\n지속적으로 베스트셀러 진입", MUTED_GREEN),
        ("AI 에이전트 개발", "LLM 기반 에이전트 개발 도서가\n새로운 카테고리로 부상", CORAL),
    ]

    for i, (title, desc, accent) in enumerate(trends):
        y = Inches(2.7) + i * Inches(1.05)
        add_rect(slide, Inches(7.5), y, Inches(0.06), Inches(0.8), accent)
        add_text_box(slide, Inches(7.8), y, Inches(4.4), Inches(0.3),
                     title, font_size=12, color=accent, bold=True)
        add_multiline_text(slide, Inches(7.8), y + Inches(0.3), Inches(4.4), Inches(0.5),
                           desc.split('\n'), font_size=10, color=SLATE_MID, line_spacing=1.4)

    set_notes(slide, "AI/LLM 카테고리 트렌드를 분석하면, Claude Code 관련 도서가 폭발적으로 "
              "성장하고 있습니다. 코워크, 하네스 등 Claude 생태계가 확장되면서 관련 도서가 "
              "급증하고 있으며, Gemini 활용서도 높은 판매지수를 기록하고 있습니다. "
              "교육 시장에서의 AI 도입도 지속적으로 이루어지고 있으며, AI 에이전트 개발이 "
              "새로운 카테고리로 부상하고 있습니다.")


def slide_08_category_trend_coding(prs, books):
    """Slide 8: Coding/Vibe Coding Trends"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WARM_WHITE)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), SLATE_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                 "07  카테고리 트렌드: 바이브 코딩", font_size=32, color=WHITE,
                 bold=True, font_name=FONT_TITLE)

    # Vibe coding books count
    vibe_books = [b for b in books if any(kw in b['title'].lower()
                  for kw in ['vibe', '바이브', '코딩'])]

    # Left: Vibe coding stats
    add_rect(slide, Inches(0.8), Inches(1.8), Inches(5.8), Inches(2.5), WHITE)
    add_text_box(slide, Inches(1.1), Inches(1.9), Inches(5), Inches(0.5),
                 "바이브 코딩 시장 현황", font_size=16, color=DEEP_TEAL,
                 bold=True, font_name=FONT_TITLE)

    vibe_stats = [
        ("관련 도서 수", f"{len(vibe_books)}권"),
        ("평균 판매지수", f"{sum(b['sales_index'] for b in vibe_books) // len(vibe_books):,}" if vibe_books else "0"),
        ("주요 플랫폼", "Cursor, Claude Code, Codex"),
        ("타겟 독자", "코딩 비전공자, 1인 창업자"),
    ]

    for i, (label, value) in enumerate(vibe_stats):
        x = Inches(1.1 + i * 1.4)
        add_text_box(slide, x, Inches(2.5), Inches(1.3), Inches(0.3),
                     label, font_size=10, color=SLATE_MID)
        add_text_box(slide, x, Inches(2.8), Inches(1.3), Inches(0.4),
                     value, font_size=14, color=DEEP_TEAL, bold=True)

    # Right: Key vibe coding books
    add_rect(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(2.5), WHITE)
    add_text_box(slide, Inches(7.3), Inches(1.9), Inches(4.9), Inches(0.5),
                 "주요 바이브 코딩 도서", font_size=16, color=WOOD_BROWN,
                 bold=True, font_name=FONT_TITLE)

    top_vibe = sorted(vibe_books, key=lambda x: x['sales_index'], reverse=True)[:4]
    for i, b in enumerate(top_vibe):
        y = Inches(2.5) + i * Inches(0.45)
        add_circle(slide, Inches(7.3), y + Inches(0.05), Inches(0.25), DEEP_TEAL)
        add_text_box(slide, Inches(7.3), y + Inches(0.05), Inches(0.25), Inches(0.25),
                     str(i+1), font_size=9, color=WHITE, bold=True,
                     align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        title = b['title'][:35] + '...' if len(b['title']) > 35 else b['title']
        add_text_box(slide, Inches(7.7), y, Inches(4.5), Inches(0.35),
                     title, font_size=10, color=SLATE_DARK)

    # Bottom: Opportunity insight
    add_rect(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.5), SLATE_DARK)
    add_text_box(slide, Inches(1.2), Inches(4.8), Inches(4), Inches(0.4),
                 "기회 분석", font_size=16, color=WOOD_BROWN, bold=True, font_name=FONT_TITLE)

    opportunities = [
        ("1인 창업 시장", "바이브 코딩으로 앱을 만들어 수익화하는\n트렌드가 확산 중", DEEP_TEAL),
        ("비개발자 타겟", "코딩을 모르는 교사, 직장인, 사업자 등\n넓은 타겟층 확보 가능", WOOD_BROWN),
        ("수익화 가이드", "앱 마켓 등록, 인앱 결제 등\n실전 수익화 내용 포함 도서 수요 증가", MUTED_GREEN),
        ("도구 다양화", "Cursor, Claude Code, Codex 등\n다양한 도구별 가이드 필요", CORAL),
    ]

    for i, (title, desc, accent) in enumerate(opportunities):
        x = Inches(1.2 + i * 2.9)
        add_text_box(slide, x, Inches(5.3), Inches(2.6), Inches(0.35),
                     title, font_size=12, color=accent, bold=True)
        add_multiline_text(slide, x, Inches(5.7), Inches(2.6), Inches(1.0),
                           desc.split('\n'), font_size=10, color=LIGHT_GRAY, line_spacing=1.5)

    set_notes(slide, "바이브 코딩 트렌드를 살펴보면, 코딩을 모르는 일반인도 앱을 만들 수 있는 "
              "이 새로운 방법론이 폭발적인 인기를 끌고 있습니다. Cursor, Claude Code, Codex 등 "
              "다양한 도구가 등장했으며, 1인 창업 시장에서의 수요가 급증하고 있습니다. "
              "이 시장은 비개발자 타겟층이 넓어 잠재력이 크며, 실전 수익화 가이드 수요도 "
              "증가하고 있습니다.")


def slide_09_category_trend_edu(prs, books):
    """Slide 9: Education/EdTech Trends"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WARM_WHITE)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), SLATE_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                 "08  카테고리 트렌드: 교육/에듀테크", font_size=32, color=WHITE,
                 bold=True, font_name=FONT_TITLE)

    # Education books
    edu_books = [b for b in books if any(kw in b['title'].lower()
                 for kw in ['교육', '교사', '수업', '에듀테크', '학교'])]

    # Left: Education market overview
    add_rect(slide, Inches(0.8), Inches(1.8), Inches(5.8), Inches(2.2), WHITE)
    add_text_box(slide, Inches(1.1), Inches(1.9), Inches(5), Inches(0.5),
                 "교육 시장 현황", font_size=16, color=DEEP_TEAL,
                 bold=True, font_name=FONT_TITLE)

    edu_stats = [
        ("교육 도서 수", f"{len(edu_books)}권"),
        ("평균 판매지수", f"{sum(b['sales_index'] for b in edu_books) // len(edu_books):,}" if edu_books else "0"),
        ("최고 판매지수", f"{max(b['sales_index'] for b in edu_books):,}" if edu_books else "0"),
    ]

    for i, (label, value) in enumerate(edu_stats):
        x = Inches(1.1 + i * 1.8)
        add_text_box(slide, x, Inches(2.5), Inches(1.6), Inches(0.3),
                     label, font_size=11, color=SLATE_MID)
        add_text_box(slide, x, Inches(2.8), Inches(1.6), Inches(0.4),
                     value, font_size=18, color=DEEP_TEAL, bold=True, font_name=FONT_TITLE)

    # Right: Education tools
    add_rect(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(2.2), WHITE)
    add_text_box(slide, Inches(7.3), Inches(1.9), Inches(4.9), Inches(0.5),
                 "주요 에듀테크 도구", font_size=16, color=WOOD_BROWN,
                 bold=True, font_name=FONT_TITLE)

    tools = [
        ("ChatGPT / Gemini", "AI 기반 수업 설계 및 학생 질문 대응"),
        ("Canva", "수업 자료 및 발표 자료 제작"),
        ("Notion", "수업 관리 및 학생 포트폴리오"),
        ("Padlet", "협업 및 토론 활동"),
    ]

    for i, (tool, desc) in enumerate(tools):
        y = Inches(2.5) + i * Inches(0.38)
        add_text_box(slide, Inches(7.3), y, Inches(2.0), Inches(0.3),
                     tool, font_size=11, color=DEEP_TEAL, bold=True)
        add_text_box(slide, Inches(9.3), y, Inches(3.0), Inches(0.3),
                     desc, font_size=10, color=SLATE_MID)

    # Bottom: Education opportunities
    add_rect(slide, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.8), WHITE)
    add_text_box(slide, Inches(1.1), Inches(4.4), Inches(5), Inches(0.5),
                 "교육 시장 기회 요소", font_size=16, color=MUTED_GREEN,
                 bold=True, font_name=FONT_TITLE)

    opps = [
        ("2022 개정 교육과정 연계", "개정 교육과정에 맞춘 AI 수업 활용 가이드가\n지속적으로 인기", DEEP_TEAL),
        ("과목별 특화", "영어, 수학, 과학 등 과목별 에듀테크 도서\n수요 확대", WOOD_BROWN),
        ("교사 업무 경감", "AI를 활용한 업무 자동화, 채점, 생기부 작성\n등 실전 가이드 수요", MUTED_GREEN),
        ("학교 도서관 시장", "학교 도서관 납품을 위한 도서 수요\n안정적 시장 확보", CORAL),
    ]

    for i, (title, desc, accent) in enumerate(opps):
        x = Inches(1.1 + i * 2.9)
        y = Inches(5.0)
        add_rect(slide, x, y, Inches(2.6), Inches(1.8), CREAM)
        add_rect(slide, x, y, Inches(2.6), Inches(0.06), accent)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.2), Inches(0.35),
                     title, font_size=12, color=accent, bold=True)
        add_multiline_text(slide, x + Inches(0.2), y + Inches(0.6), Inches(2.2), Inches(1.0),
                           desc.split('\n'), font_size=10, color=SLATE_MID, line_spacing=1.5)

    set_notes(slide, "교육/에듀테크 시장을 분석하면, 교사 대상 AI 활용 도서가 지속적으로 "
              "인기를 끌고 있습니다. 2022 개정 교육과정에 맞춘 도서, 과목별 특화 도서, "
              "교사 업무 경감을 위한 실전 가이드 등 다양한 수요가 있습니다. "
              "학교 도서관 시장은 안정적인 납품 수요가 있어 신규 도서 기획 시 "
              "고려할 만한 요소입니다.")


def slide_10_audience(prs, books):
    """Slide 10: Target Audience Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WARM_WHITE)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), SLATE_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                 "09  타겟 독자층 분석", font_size=32, color=WHITE,
                 bold=True, font_name=FONT_TITLE)

    # Audience segments
    audiences = [
        ("일반인/비개발자", "코딩을 모르지만 AI 도구로\n무언가를 만들고 싶은 사람들",
         "바이브 코딩, ChatGPT 활용서", "가장 넓은 타겟층", DEEP_TEAL),
        ("직장인", "업무 효율을 높이려는\n직장인 및 관리자",
         "AI 업무 활용법, 자동화 가이드", "높은 구매력", WOOD_BROWN),
        ("교사/교육자", "AI를 수업에 도입하려는\n초중고 교사",
         "에듀테크 수업 활용 가이드", "안정적 수요", MUTED_GREEN),
        ("개발자/엔지니어", "AI 기반 개발 도구를\n학습하는 개발자",
         "AI 에이전트, Claude Code 서적", "전문성 중시", CORAL),
    ]

    for i, (title, desc, books_example, note, accent) in enumerate(audiences):
        x = Inches(0.8 + i * 3.1)
        y = Inches(1.8)

        add_rect(slide, x, y, Inches(2.8), Inches(4.5), WHITE)
        add_rect(slide, x, y, Inches(2.8), Inches(0.8), accent)

        add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(2.4), Inches(0.5),
                     title, font_size=16, color=WHITE, bold=True,
                     align=PP_ALIGN.CENTER, font_name=FONT_TITLE)

        add_multiline_text(slide, x + Inches(0.2), y + Inches(1.0), Inches(2.4), Inches(1.0),
                           desc.split('\n'), font_size=11, color=SLATE_MID,
                           line_spacing=1.5, align=PP_ALIGN.CENTER)

        add_rect(slide, x + Inches(0.2), y + Inches(2.2), Inches(2.4), Inches(0.02), LIGHT_GRAY)

        add_text_box(slide, x + Inches(0.2), y + Inches(2.4), Inches(2.4), Inches(0.3),
                     "선호 도서 유형", font_size=10, color=SLATE_MID,
                     align=PP_ALIGN.CENTER, italic=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(2.7), Inches(2.4), Inches(0.5),
                     books_example, font_size=10, color=DEEP_TEAL,
                     align=PP_ALIGN.CENTER, bold=True)

        add_rect(slide, x + Inches(0.2), y + Inches(3.4), Inches(2.4), Inches(0.02), LIGHT_GRAY)

        # Note badge
        badge = add_rounded_rect(slide, x + Inches(0.5), y + Inches(3.6), Inches(1.8), Inches(0.4), accent)
        add_text_box(slide, x + Inches(0.5), y + Inches(3.6), Inches(1.8), Inches(0.4),
                     note, font_size=10, color=WHITE, bold=True,
                     align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # Bottom insight
    add_rect(slide, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.6), SLATE_DARK)
    add_text_box(slide, Inches(1.2), Inches(6.65), Inches(11), Inches(0.5),
                 "핵심 인사이트: 비개발자 타겟층이 가장 넓으며, AI 도구를 활용한 실전 가이드가 "
                 "가장 높은 수요를 보이고 있습니다.",
                 font_size=12, color=LIGHT_GRAY)

    set_notes(slide, "타겟 독자층을 4가지로 분류하면, 일반인/비개발자가 가장 넓은 타겟층이며, "
              "바이브 코딩과 ChatGPT 활용서에 대한 수요가 높습니다. 직장인은 업무 효율화를 위한 "
              "AI 활용서를 선호하며, 교사는 안정적인 수요를 보유하고 있습니다. "
              "개발자는 전문성을 중시하는 서적을 선호합니다.")


def slide_11_book_concept(prs):
    """Slide 11: New Book Concept"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, SLATE_DARK)

    # Left decorative bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, WOOD_BROWN)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                 "10  신규 도서 기획 컨셉", font_size=32, color=WHITE,
                 bold=True, font_name=FONT_TITLE)

    # Main concept card
    add_rect(slide, Inches(0.8), Inches(1.4), Inches(7.0), Inches(5.5), SLATE_MID)
    add_rect(slide, Inches(0.8), Inches(1.4), Inches(7.0), Inches(0.08), WOOD_BROWN)

    add_text_box(slide, Inches(1.2), Inches(1.7), Inches(6.2), Inches(0.5),
                 "제안 도서명", font_size=14, color=WOOD_BROWN, bold=True)
    add_text_box(slide, Inches(1.2), Inches(2.2), Inches(6.2), Inches(0.8),
                 "된다! AI 바이브 코딩 완전 정복\nwith Claude Code & Gemini",
                 font_size=28, color=WHITE, bold=True, font_name=FONT_TITLE)

    add_text_box(slide, Inches(1.2), Inches(3.2), Inches(6.2), Inches(0.3),
                 "코딩 몰라도 OK! AI와 대화로 나만의 앱을 만드는 실전 가이드",
                 font_size=14, color=SOFT_BLUE, italic=True)

    add_rect(slide, Inches(1.2), Inches(3.7), Inches(5.5), Inches(0.02), LIGHT_GRAY)

    # Book features
    features = [
        ("대상 독자", "코딩 비전공자, 1인 창업자, 직장인, 교사"),
        ("분량", "350페이지 내외"),
        ("정가", "27,000원 (판매가 24,300원)"),
        ("출간 목표", "2026년 Q4"),
        ("차별점", "멀티 도구 비교 + 실전 수익화 가이드"),
    ]

    for i, (label, value) in enumerate(features):
        y = Inches(3.9) + i * Inches(0.55)
        add_text_box(slide, Inches(1.2), y, Inches(1.8), Inches(0.4),
                     label, font_size=12, color=WOOD_BROWN, bold=True)
        add_text_box(slide, Inches(3.0), y, Inches(4.4), Inches(0.4),
                     value, font_size=12, color=LIGHT_GRAY)

    # Right side: Key selling points
    add_rect(slide, Inches(8.2), Inches(1.4), Inches(4.3), Inches(5.5), WHITE)

    add_text_box(slide, Inches(8.5), Inches(1.6), Inches(3.7), Inches(0.5),
                 "핵심 판매 포인트", font_size=16, color=DEEP_TEAL,
                 bold=True, font_name=FONT_TITLE)

    selling_points = [
        ("멀티 도구 비교", "Cursor, Claude Code, Gemini, Codex 등\n주요 도구를 한 권으로 비교", DEEP_TEAL),
        ("단계별 프로젝트", "실제 앱을 만들면서 배우는\n Hands-on 학습 방식", WOOD_BROWN),
        ("수익화 가이드", "앱 마켓 등록부터 인앱 결제,\n구독 상품 설정까지", MUTED_GREEN),
        ("비개발자 친화", "코딩 개념 설명 없이\nAI 대화 중심으로 구성", CORAL),
        ("최신 트렌드 반영", "2026년 최신 AI 도구 및\n업데이트 내용 포함", DARK_TEAL),
    ]

    for i, (title, desc, accent) in enumerate(selling_points):
        y = Inches(2.3) + i * Inches(0.95)
        add_circle(slide, Inches(8.5), y + Inches(0.05), Inches(0.35), accent)
        add_text_box(slide, Inches(8.5), y + Inches(0.05), Inches(0.35), Inches(0.35),
                     str(i+1), font_size=11, color=WHITE, bold=True,
                     align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(9.0), y, Inches(3.3), Inches(0.3),
                     title, font_size=12, color=accent, bold=True)
        add_multiline_text(slide, Inches(9.0), y + Inches(0.3), Inches(3.3), Inches(0.5),
                           desc.split('\n'), font_size=10, color=SLATE_MID, line_spacing=1.3)

    set_notes(slide, "신규 도서 기획 컨셉입니다. '된다! AI 바이브 코딩 완전 정복 with "
              "Claude Code & Gemini'로 제안하며, 코딩을 모르는 일반인도 AI와 대화로 "
              "나만의 앱을 만들 수 있는 실전 가이드를 목표로 합니다. 멀티 도구 비교, "
              "단계별 프로젝트, 수익화 가이드 등 기존 도서와 차별화된 포인트를 "
              "제시합니다.")


def slide_12_competitive(prs, books):
    """Slide 12: Competitive Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WARM_WHITE)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), SLATE_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                 "11  경쟁 도서 분석", font_size=32, color=WHITE,
                 bold=True, font_name=FONT_TITLE)

    # Competitor comparison table
    competitors = [
        ("요즘 바이브 코딩\n클로드 코드 완벽 가이드", "골든래빗", "21,600원", "28,542", "단일 도구\n(Claude Code)"),
        ("클로드 코드 마스터", "한빛미디어", "32,400원", "14,175", "단일 도구\n(Claude Code)"),
        ("바이브 코딩 with\n커서 AI", "골든래빗", "27,000원", "3,414", "단일 도구\n(Cursor)"),
        ("된다! 하루 만에\n끝내는 바이브 코딩", "이지스퍼블리싱", "19,800원", "21,051", "입문서\n(범용)"),
    ]

    # Table header
    header_y = Inches(1.6)
    add_rect(slide, Inches(0.8), header_y, Inches(11.7), Inches(0.6), SLATE_MID)

    cols = [("경쟁 도서", 0.8, 3.0), ("출판사", 3.8, 1.8), ("가격", 5.6, 1.2),
            ("판매지수", 6.8, 1.5), ("특징", 8.3, 3.0)]

    for text, x, w in cols:
        add_text_box(slide, Inches(x), header_y + Inches(0.12), Inches(w), Inches(0.35),
                     text, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Table rows
    for i, (title, pub, price, index, feature) in enumerate(competitors):
        row_y = Inches(2.3 + i * 0.9)
        bg_color = WHITE if i % 2 == 0 else CREAM
        add_rect(slide, Inches(0.8), row_y, Inches(11.7), Inches(0.85), bg_color)

        lines = title.split('\n')
        add_multiline_text(slide, Inches(1.0), row_y + Inches(0.1), Inches(2.8), Inches(0.65),
                           lines, font_size=11, color=SLATE_DARK, bold=True,
                           line_spacing=1.2, align=PP_ALIGN.LEFT)

        add_text_box(slide, Inches(3.8), row_y + Inches(0.2), Inches(1.8), Inches(0.45),
                     pub, font_size=11, color=SLATE_MID, align=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(5.6), row_y + Inches(0.2), Inches(1.2), Inches(0.45),
                     price, font_size=11, color=SLATE_DARK, align=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(6.8), row_y + Inches(0.2), Inches(1.5), Inches(0.45),
                     index, font_size=12, color=DEEP_TEAL, bold=True, align=PP_ALIGN.CENTER)

        feat_lines = feature.split('\n')
        add_multiline_text(slide, Inches(8.5), row_y + Inches(0.1), Inches(2.8), Inches(0.65),
                           feat_lines, font_size=10, color=SLATE_MID,
                           line_spacing=1.2, align=PP_ALIGN.CENTER)

    # Our differentiation
    add_rect(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(1.2), SLATE_DARK)
    add_text_box(slide, Inches(1.2), Inches(6.1), Inches(3), Inches(0.4),
                 "우리의 차별화 포인트", font_size=14, color=WOOD_BROWN, bold=True)

    diffs = [
        "멀티 도구 비교 (Cursor + Claude Code + Gemini + Codex)",
        "실전 수익화 가이드 포함",
        "비개발자 최적화 (코딩 개념 설명 최소화)",
    ]

    for i, diff in enumerate(diffs):
        x = Inches(1.2 + i * 3.8)
        add_circle(slide, x, Inches(6.55), Inches(0.2), WOOD_BROWN)
        add_text_box(slide, x + Inches(0.3), Inches(6.5), Inches(3.3), Inches(0.5),
                     diff, font_size=11, color=LIGHT_GRAY)

    set_notes(slide, "경쟁 도서를 분석하면, 대부분 단일 도구(Claude Code 또는 Cursor)에 "
              "집중하고 있습니다. 우리의 차별화 포인트는 멀티 도구 비교, 실전 수익화 가이드, "
              "비개발자 최적화입니다. 기존 도서들이 놓치고 있는 부분을 공략하여 "
              "시장에서의 경쟁 우위를 확보하겠습니다.")


def slide_13_marketing(prs):
    """Slide 13: Marketing Strategy"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WARM_WHITE)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), SLATE_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                 "12  마케팅 전략", font_size=32, color=WHITE,
                 bold=True, font_name=FONT_TITLE)

    # Marketing channels
    channels = [
        ("온라인 마케팅", [
            "유튜브 바이브 코딩 채널 협업",
            "블로그/카페 체험단 운영",
            "인스타그램 릴스 콘텐츠 마케팅",
            "검색광고 (네이버, 구글)",
        ], DEEP_TEAL),
        ("오프라인 마케팅", [
            "서점 진열 및 전면 배치",
            "교육기관 세미나/웨비나",
            "IT 컨퍼런스 부스 참여",
            "서점 사인회 이벤트",
        ], WOOD_BROWN),
        ("콘텐츠 마케팅", [
            "무료 챕터 공개로 관심 유도",
            "바이브 코딩 챌린지 이벤트",
            "독자 후기 공유 캠페인",
            "AI 도구별 비교 인포그래픽",
        ], MUTED_GREEN),
    ]

    for i, (title, items, accent) in enumerate(channels):
        x = Inches(0.8 + i * 4.1)
        y = Inches(1.8)

        add_rect(slide, x, y, Inches(3.7), Inches(3.5), WHITE)
        add_rect(slide, x, y, Inches(3.7), Inches(0.7), accent)

        add_text_box(slide, x + Inches(0.3), y + Inches(0.1), Inches(3.1), Inches(0.5),
                     title, font_size=16, color=WHITE, bold=True,
                     align=PP_ALIGN.CENTER, font_name=FONT_TITLE)

        for j, item in enumerate(items):
            item_y = y + Inches(0.9) + j * Inches(0.6)
            add_circle(slide, x + Inches(0.3), item_y + Inches(0.08), Inches(0.2), accent)
            add_text_box(slide, x + Inches(0.3), item_y + Inches(0.08), Inches(0.2), Inches(0.2),
                         str(j+1), font_size=9, color=WHITE, bold=True,
                         align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            add_text_box(slide, x + Inches(0.65), item_y, Inches(2.8), Inches(0.4),
                         item, font_size=11, color=SLATE_DARK)

    # Bottom: Key metrics
    add_rect(slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.5), SLATE_DARK)

    metrics = [
        ("출간 전 사전예약", "출간 2개월 전\n온라인 사전예약 시작"),
        ("출간 첫 주", "온·오프라인\n집중 마케팅"),
        ("출간 후 1개월", "독자 리뷰 확보\n판매지수 관리"),
        ("출간 후 3개월", "개정판 검토\n후속작 기획"),
    ]

    for i, (title, desc) in enumerate(metrics):
        x = Inches(1.2 + i * 2.9)
        add_text_box(slide, x, Inches(5.7), Inches(2.5), Inches(0.3),
                     title, font_size=12, color=WOOD_BROWN, bold=True,
                     align=PP_ALIGN.CENTER)
        add_multiline_text(slide, x, Inches(6.1), Inches(2.5), Inches(0.7),
                           desc.split('\n'), font_size=11, color=LIGHT_GRAY,
                           line_spacing=1.3, align=PP_ALIGN.CENTER)

        if i < len(metrics) - 1:
            add_text_box(slide, Inches(x.inches + 2.5), Inches(5.9), Inches(0.3), Inches(0.4),
                         "→", font_size=16, color=WOOD_BROWN, align=PP_ALIGN.CENTER,
                         valign=MSO_ANCHOR.MIDDLE)

    set_notes(slide, "마케팅 전략은 온라인, 오프라인, 콘텐츠 마케팅 3가지 축으로 구성됩니다. "
              "온라인은 유튜브 바이브 코딩 채널 협업과 블로그 체험단을, 오프라인은 서점 진열과 "
              "교육기관 세미나를, 콘텐츠는 무료 챕터 공개와 바이브 코딩 챌린지 이벤트를 "
              "추진하겠습니다. 출간 전 사전예약부터 출간 후 3개월까지 체계적인 마케팅 "
              "로드맵을 수립했습니다.")


def slide_14_timeline(prs):
    """Slide 14: Timeline & Milestones"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WARM_WHITE)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), SLATE_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                 "13  일정 및 마일스톤", font_size=32, color=WHITE,
                 bold=True, font_name=FONT_TITLE)

    # Timeline
    phases = [
        ("1단계", "기획 및 리서치", "2026.08 ~ 2026.09", "2개월",
         ["시장 조사 완료", "도서 구성 확정", "저자 섭외"], DEEP_TEAL),
        ("2단계", "원고 집필", "2026.10 ~ 2027.01", "4개월",
         ["챕터별 원고 집필", "편집 및 교정", "일러스트 제작"], WOOD_BROWN),
        ("3단계", "편집 및 디자인", "2027.02 ~ 2027.03", "2개월",
         ["레이아웃 디자인", "표지 디자인", "인쇄 준비"], MUTED_GREEN),
        ("4단계", "출간 및 마케팅", "2027.04", "1개월",
         ["출간", "마케팅 캠페인", "서점 배포"], CORAL),
    ]

    # Timeline line
    add_rect(slide, Inches(1.5), Inches(2.3), Inches(10.3), Inches(0.04), LIGHT_GRAY)

    for i, (phase, title, period, duration, items, accent) in enumerate(phases):
        x = Inches(1.2 + i * 2.8)

        # Timeline dot
        add_circle(slide, x + Inches(1.0), Inches(2.1), Inches(0.4), accent)
        add_text_box(slide, x + Inches(1.0), Inches(2.1), Inches(0.4), Inches(0.4),
                     str(i+1), font_size=12, color=WHITE, bold=True,
                     align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

        # Phase card
        card_y = Inches(2.8)
        add_rect(slide, x, card_y, Inches(2.5), Inches(3.8), WHITE)
        add_rect(slide, x, card_y, Inches(2.5), Inches(0.06), accent)

        add_text_box(slide, x + Inches(0.15), card_y + Inches(0.2), Inches(2.2), Inches(0.3),
                     phase, font_size=10, color=accent, bold=True)
        add_text_box(slide, x + Inches(0.15), card_y + Inches(0.5), Inches(2.2), Inches(0.4),
                     title, font_size=14, color=SLATE_DARK, bold=True, font_name=FONT_TITLE)
        add_text_box(slide, x + Inches(0.15), card_y + Inches(1.0), Inches(2.2), Inches(0.3),
                     period, font_size=10, color=SLATE_MID)
        add_text_box(slide, x + Inches(0.15), card_y + Inches(1.3), Inches(2.2), Inches(0.3),
                     duration, font_size=11, color=accent, bold=True)

        add_rect(slide, x + Inches(0.15), card_y + Inches(1.7), Inches(2.2), Inches(0.02), LIGHT_GRAY)

        for j, item in enumerate(items):
            item_y = card_y + Inches(1.9) + j * Inches(0.5)
            add_circle(slide, x + Inches(0.15), item_y + Inches(0.08), Inches(0.15), accent)
            add_text_box(slide, x + Inches(0.4), item_y, Inches(1.9), Inches(0.35),
                         item, font_size=10, color=SLATE_MID)

    # Bottom summary
    add_rect(slide, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.5), SLATE_DARK)
    add_text_box(slide, Inches(1.2), Inches(6.85), Inches(11), Inches(0.4),
                 "총 개발 기간: 약 9개월  |  목표 출간일: 2027년 4월  |  "
                 "사전 마케팅 시작: 2027년 2월",
                 font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    set_notes(slide, "일정 및 마일스톤입니다. 총 9개월의 개발 기간을 거쳐 2027년 4월 출간을 "
              "목표로 합니다. 1단계 기획 및 리서치 2개월, 2단계 원고 집필 4개월, "
              "3단계 편집 및 디자인 2개월, 4단계 출간 및 마케팅 1개월로 구성됩니다. "
              "출간 2개월 전부터 사전 마케팅을 시작하여 초기 판매를 활성화하겠습니다.")


def slide_15_closing(prs):
    """Slide 15: Closing/Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, SLATE_DARK)

    # Decorative elements
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, WOOD_BROWN)

    # Decorative circles
    add_circle(slide, Inches(9.5), Inches(0.5), Inches(2.5), SLATE_MID)
    add_circle(slide, Inches(10.5), Inches(2.5), Inches(1.5), DEEP_TEAL)
    add_circle(slide, Inches(0.5), Inches(5.5), Inches(1.8), SLATE_MID)

    # Main text
    add_text_box(slide, Inches(1.5), Inches(1.5), Inches(8), Inches(1.0),
                 "감사합니다", font_size=48, color=WHITE,
                 bold=True, font_name=FONT_TITLE, align=PP_ALIGN.LEFT)

    add_rect(slide, Inches(1.5), Inches(2.7), Inches(3), Inches(0.04), WOOD_BROWN)

    add_text_box(slide, Inches(1.5), Inches(3.1), Inches(8), Inches(0.6),
                 "신규 도서 기획안 검토 및 피드백 부탁드립니다",
                 font_size=18, color=SOFT_BLUE, font_name=FONT_BODY)

    # Contact info
    add_rect(slide, Inches(1.5), Inches(4.2), Inches(5), Inches(2.0), SLATE_MID)
    add_text_box(slide, Inches(1.8), Inches(4.4), Inches(4.4), Inches(0.4),
                 "연락처", font_size=14, color=WOOD_BROWN, bold=True)

    contacts = [
        "ABC Publishing  |  IT & Technology Division",
        "이메일: contact@abc-publishing.co.kr",
        "전화: 02-XXX-XXXX",
    ]

    for i, contact in enumerate(contacts):
        add_text_box(slide, Inches(1.8), Inches(4.9) + i * Inches(0.4), Inches(4.4), Inches(0.35),
                     contact, font_size=12, color=LIGHT_GRAY)

    # Bottom bar
    add_rect(slide, Inches(0), Inches(6.8), SLIDE_W, Inches(0.7), SLATE_MID)
    add_text_box(slide, Inches(1.5), Inches(6.85), Inches(10), Inches(0.5),
                 "ABC Publishing  |  \"books that inspire\"",
                 font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

    set_notes(slide, "감사합니다. 오늘 제안드린 신규 도서 기획안에 대한 검토 및 "
              "피드백을 부탁드립니다. 추가적인 질문이나 수정 사항이 있으시면 "
              "언제든지 연락 주시기 바랍니다. 감사합니다.")


# ── Main Execution ──

def main():
    # Load data
    csv_path = os.path.join(os.path.dirname(__file__), 'yes24_it_mobile_bestsellers.csv')
    books = load_data(csv_path)
    print(f"Loaded {len(books)} books from CSV")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Generate slides
    print("Generating slides...")
    slide_01_title(prs)
    print("  Slide 01: Title")
    slide_02_overview(prs)
    print("  Slide 02: Overview")
    slide_03_market_overview(prs)
    print("  Slide 03: Market Overview")
    slide_04_top10(prs, books)
    print("  Slide 04: Top 10")
    slide_05_publisher_analysis(prs, books)
    print("  Slide 05: Publisher Analysis")
    slide_06_price_analysis(prs, books)
    print("  Slide 06: Price Analysis")
    slide_07_category_trend_ai(prs, books)
    print("  Slide 07: AI Trends")
    slide_08_category_trend_coding(prs, books)
    print("  Slide 08: Coding Trends")
    slide_09_category_trend_edu(prs, books)
    print("  Slide 09: Education Trends")
    slide_10_audience(prs, books)
    print("  Slide 10: Target Audience")
    slide_11_book_concept(prs)
    print("  Slide 11: Book Concept")
    slide_12_competitive(prs, books)
    print("  Slide 12: Competitive Analysis")
    slide_13_marketing(prs)
    print("  Slide 13: Marketing Strategy")
    slide_14_timeline(prs)
    print("  Slide 14: Timeline")
    slide_15_closing(prs)
    print("  Slide 15: Closing")

    # Save
    output_path = os.path.join(os.path.dirname(__file__), 'nordic_book_proposal.pptx')
    prs.save(output_path)
    print(f"\nPresentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == '__main__':
    main()
